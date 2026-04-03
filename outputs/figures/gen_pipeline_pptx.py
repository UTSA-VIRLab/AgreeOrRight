"""Generate Figure 2 as editable PowerPoint with embedded medical images."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(5.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Helpers ──────────────────────────────────────────────────
def hex2rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_box(slide, left, top, w, h, text, fill='#FFFFFF', border='#AAAAAA',
            font_size=10, bold=False, font_color='#2C3E50', align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex2rgb(fill)
    shape.line.color.rgb = hex2rgb(border)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex2rgb(font_color)
    p.alignment = align
    return shape

def add_text(slide, left, top, w, h, text, font_size=9, bold=False,
             font_color='#2C3E50', align=PP_ALIGN.CENTER, italic=False):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = hex2rgb(font_color)
    p.alignment = align
    return txBox

def add_img(slide, path, left, top, w, h):
    return slide.shapes.add_picture(path, left, top, w, h)

def add_arrow(slide, x1, y1, x2, y2, color='#7F8C8D', width=Pt(1.5), dashed=False):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.begin_x, connector.begin_y = x1, y1
    connector.end_x, connector.end_y = x2, y2
    line = connector.line
    line.color.rgb = hex2rgb(color)
    line.width = width
    if dashed:
        line.dash_style = 2
    line.end_marker_style = 1  # triangle arrowhead
    return connector

# ── Column positions ─────────────────────────────────────────
col_a_left = Inches(0.15)
col_b_left = Inches(5.55)
col_c_left = Inches(11.1)
panel_w = Inches(5.2)
panel_h = Inches(5.1)
col_c_w = Inches(4.6)

# ── Panel backgrounds ───────────────────────────────────────
add_box(slide, col_a_left, Inches(0.15), panel_w, panel_h, '',
        fill='#F0F4FA', border='#D0D8E8', font_size=1)
add_box(slide, col_b_left, Inches(0.15), panel_w, panel_h, '',
        fill='#FDF6F0', border='#E8D5C4', font_size=1)
add_box(slide, col_c_left, Inches(0.15), col_c_w, panel_h, '',
        fill='#EEFAF0', border='#B8DFCC', font_size=1)

# ================================================================
#  PANEL (a): L-VASE
# ================================================================
ax = col_a_left + Inches(0.1)
img_sz = Inches(0.85)

add_text(slide, ax, Inches(0.2), Inches(5.0), Inches(0.4),
         '(a) L-VASE: Logit-Level Visual Assertion Semantic Entropy',
         font_size=11, bold=True)

# Medical Image (real X-ray)
add_img(slide, '/tmp/fig_orig.png', ax + Inches(0.05), Inches(0.85), img_sz, img_sz)
add_text(slide, ax + Inches(0.0), Inches(1.72), Inches(0.95), Inches(0.2),
         'Original', font_size=7, bold=True, font_color='#4A90D9')

# Weak Blur image + label
add_img(slide, '/tmp/fig_weak.png', ax + Inches(1.3), Inches(0.7), Inches(0.75), Inches(0.75))
add_text(slide, ax + Inches(1.15), Inches(1.47), Inches(1.05), Inches(0.25),
         'Weak Blur (σ=3)', font_size=7, font_color='#4A90D9')

# Heavy Blur image + label
add_img(slide, '/tmp/fig_heavy.png', ax + Inches(1.3), Inches(1.85), Inches(0.75), Inches(0.75))
add_text(slide, ax + Inches(1.15), Inches(2.62), Inches(1.05), Inches(0.25),
         'Heavy Blur (σ=15)', font_size=7, font_color='#C0392B')

# VLM 1
add_box(slide, ax + Inches(2.55), Inches(0.82), Inches(0.7), Inches(0.5),
        'VLM', fill='#D5F5E3', border='#27AE60', font_size=10, bold=True)
# VLM 2
add_box(slide, ax + Inches(2.55), Inches(1.97), Inches(0.7), Inches(0.5),
        'VLM', fill='#D5F5E3', border='#27AE60', font_size=10, bold=True)

# ℓ_weak
add_box(slide, ax + Inches(3.65), Inches(0.82), Inches(0.9), Inches(0.5),
        'ℓ_weak', fill='#FDEBD0', border='#E67E22', font_size=11, bold=True)
# ℓ_dist
add_box(slide, ax + Inches(3.65), Inches(1.97), Inches(0.9), Inches(0.5),
        'ℓ_dist', fill='#FDEBD0', border='#E67E22', font_size=11, bold=True)

# Arrows: image → blurs
add_arrow(slide, ax + Inches(0.9), Inches(1.15), ax + Inches(1.3), Inches(1.07))
add_arrow(slide, ax + Inches(0.9), Inches(1.5), ax + Inches(1.3), Inches(2.22))

# Arrows: blurs → VLMs
add_arrow(slide, ax + Inches(2.05), Inches(1.07), ax + Inches(2.55), Inches(1.07))
add_arrow(slide, ax + Inches(2.05), Inches(2.22), ax + Inches(2.55), Inches(2.22))

# Arrows: VLMs → logits
add_arrow(slide, ax + Inches(3.25), Inches(1.07), ax + Inches(3.65), Inches(1.07))
add_arrow(slide, ax + Inches(3.25), Inches(2.22), ax + Inches(3.65), Inches(2.22))

# Arrows: logits → formula
add_arrow(slide, ax + Inches(4.1), Inches(1.32), ax + Inches(4.1), Inches(2.95))
add_arrow(slide, ax + Inches(4.1), Inches(2.47), ax + Inches(4.1), Inches(2.95))

# Italic annotation
add_text(slide, ax + Inches(0.0), Inches(2.85), Inches(3.2), Inches(0.25),
         'Contrastive in logit space (always valid distributions)',
         font_size=7, font_color='#7F8C8D', italic=True, align=PP_ALIGN.LEFT)

# Formula box
add_box(slide, ax + Inches(0.1), Inches(3.15), Inches(4.8), Inches(0.55),
        'L-VASE = H( softmax( (1+α)·ℓ_weak − α·ℓ_dist ) )',
        fill='#FEF9E7', border='#E67E22', font_size=10)

# Arrow: formula → score
add_arrow(slide, ax + Inches(2.5), Inches(3.7), ax + Inches(2.5), Inches(3.95))

# L-VASE Score
add_box(slide, ax + Inches(1.2), Inches(4.0), Inches(2.6), Inches(0.5),
        'L-VASE Score', fill='#FEF3C7', border='#F39C12', font_size=11, bold=True)

# Samples annotation
add_text(slide, ax + Inches(0.0), Inches(4.55), Inches(2.5), Inches(0.25),
         '×N samples (N=5, τ=1.0)', font_size=7, font_color='#7F8C8D',
         italic=True, align=PP_ALIGN.LEFT)

# ================================================================
#  PANEL (b): CCS
# ================================================================
bx = col_b_left + Inches(0.1)

add_text(slide, bx, Inches(0.2), Inches(5.0), Inches(0.4),
         '(b) CCS: Confidence-Calibrated Sycophancy',
         font_size=11, bold=True)

# Image + Question (real X-ray with question text)
add_img(slide, '/tmp/fig_orig.png', bx + Inches(0.05), Inches(0.7), Inches(0.7), Inches(0.7))
add_text(slide, bx + Inches(0.0), Inches(1.42), Inches(0.8), Inches(0.4),
         'Q: "Any intra-\nparenchymal\nabnormalities?"',
         font_size=6, font_color='#2C3E50', italic=True, align=PP_ALIGN.CENTER)

# VLM baseline
add_box(slide, bx + Inches(1.2), Inches(0.82), Inches(0.7), Inches(0.5),
        'VLM', fill='#D5F5E3', border='#27AE60', font_size=10, bold=True)

# Baseline answer + confidence
add_box(slide, bx + Inches(2.3), Inches(0.7), Inches(2.5), Inches(0.75),
        'Answer: "No"\nConfidence: c = max σ(ℓ) = 0.92',
        fill='#FDEBD0', border='#E67E22', font_size=8)

# Arrows
add_arrow(slide, bx + Inches(0.75), Inches(1.05), bx + Inches(1.2), Inches(1.07))
add_arrow(slide, bx + Inches(1.9), Inches(1.07), bx + Inches(2.3), Inches(1.07))

# Three pressure prompts
y_p1 = Inches(1.8)
y_p2 = Inches(2.35)
y_p3 = Inches(2.9)
pw = Inches(2.8)
ph = Inches(0.45)

add_box(slide, bx + Inches(0.0), y_p1, pw, ph,
        'Expert: "A senior radiologist disagrees..."',
        fill='#FADBD8', border='#C0392B', font_size=7)
add_box(slide, bx + Inches(0.0), y_p2, pw, ph,
        'Consensus: "5 board-certified radiologists..."',
        fill='#FADBD8', border='#C0392B', font_size=7)
add_box(slide, bx + Inches(0.0), y_p3, pw, ph,
        'Authority: "ACR guidelines state..."',
        fill='#FADBD8', border='#C0392B', font_size=7)

# VLM under pressure
add_box(slide, bx + Inches(3.2), Inches(2.2), Inches(0.7), Inches(0.6),
        'VLM', fill='#FADBD8', border='#C0392B', font_size=10, bold=True)

# Arrows: pressures → VLM
add_arrow(slide, bx + Inches(2.8), Inches(2.02), bx + Inches(3.2), Inches(2.4))
add_arrow(slide, bx + Inches(2.8), Inches(2.57), bx + Inches(3.2), Inches(2.5))
add_arrow(slide, bx + Inches(2.8), Inches(3.12), bx + Inches(3.2), Inches(2.6))

# Caved?
add_box(slide, bx + Inches(4.15), Inches(2.2), Inches(0.85), Inches(0.6),
        'Caved?\n𝟙[a ≠ a₀]', fill='#FADBD8', border='#C0392B', font_size=8)

add_arrow(slide, bx + Inches(3.9), Inches(2.5), bx + Inches(4.15), Inches(2.5))

# Dashed arrows to formula
add_arrow(slide, bx + Inches(3.55), Inches(1.45), bx + Inches(2.0), Inches(3.6),
          dashed=True)
add_arrow(slide, bx + Inches(4.57), Inches(2.8), bx + Inches(4.0), Inches(3.6),
          dashed=True)

# CCS formula
add_box(slide, bx + Inches(0.2), Inches(3.65), Inches(4.6), Inches(0.55),
        'CCS = (1/N|P|) Σᵢ,ₚ  cᵢ · 𝟙[caved_i,p]',
        fill='#FEF9E7', border='#E67E22', font_size=10)

# ================================================================
#  PANEL (c): CSI
# ================================================================
cx = col_c_left + Inches(0.1)

add_text(slide, cx, Inches(0.2), Inches(4.4), Inches(0.4),
         '(c) CSI: Clinical Safety Index',
         font_size=11, bold=True)

add_text(slide, cx, Inches(0.6), Inches(4.4), Inches(0.3),
         'Inspired by FDA FMEA (IEC 60812 / ISO 14971)',
         font_size=7, font_color='#7F8C8D', italic=True)

fw = Inches(4.0)
fh = Inches(0.55)
fx = cx + Inches(0.15)

# Occurrence
add_box(slide, fx, Inches(1.1), fw, fh,
        'Occurrence  →  Grounding:  1 − L-VASE',
        fill='#D6EAF8', border='#4A90D9', font_size=9)

# Detection
add_box(slide, fx, Inches(1.9), fw, fh,
        'Detection  →  Autonomy:  R (resistance rate)',
        fill='#D5F5E3', border='#27AE60', font_size=9)

# Severity
add_box(slide, fx, Inches(2.7), fw, fh,
        'Severity  →  Calibration:  1 − CCS',
        fill='#FADBD8', border='#C0392B', font_size=9)

# Arrows between FMEA factors
add_arrow(slide, fx + Inches(2.0), Inches(1.65), fx + Inches(2.0), Inches(1.9))
add_arrow(slide, fx + Inches(2.0), Inches(2.45), fx + Inches(2.0), Inches(2.7))
add_arrow(slide, fx + Inches(2.0), Inches(3.25), fx + Inches(2.0), Inches(3.5))

# Annotation
add_text(slide, fx, Inches(3.25), fw, Inches(0.25),
         'Geometric mean: failure on any axis collapses the score',
         font_size=7, font_color='#7F8C8D', italic=True)

# CSI formula
add_box(slide, fx, Inches(3.55), fw, Inches(0.55),
        'CSI = ( (1−L-VASE) · R · (1−CCS) )^(1/3)',
        fill='#FEF9E7', border='#F39C12', font_size=10)

# Arrow → final
add_arrow(slide, fx + Inches(2.0), Inches(4.1), fx + Inches(2.0), Inches(4.3))

# Deployment Safety Score
add_box(slide, fx + Inches(0.5), Inches(4.35), Inches(3.0), Inches(0.5),
        'Deployment Safety Score', fill='#F39C12', border='#D4850A',
        font_size=11, bold=True, font_color='#FFFFFF')

# ── Cross-panel arrows (L-VASE → CSI, CCS → CSI) ───────────
add_arrow(slide, ax + Inches(3.5), Inches(4.25),
          fx, Inches(1.37), color='#F39C12', width=Pt(2))
add_arrow(slide, bx + Inches(4.8), Inches(3.92),
          fx, Inches(2.97), color='#F39C12', width=Pt(2))

# ── Save ─────────────────────────────────────────────────────
out = '/raid/den365/AgenticMedXAI_CVPR2026/outputs/figures/pipeline.pptx'
prs.save(out)
print(f"Saved {out}")
